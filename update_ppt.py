import warnings
from urllib3.exceptions import InsecureRequestWarning

warnings.simplefilter('ignore', InsecureRequestWarning)

import os
import sys
import requests
from msal import ConfidentialClientApplication
from openpyxl import load_workbook
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from dotenv import load_dotenv

load_dotenv()

config = {
    "client_id": os.getenv("CLIENT_ID"),
    "authority": os.getenv("AUTHORITY"),
    "client_secret": os.getenv("CLIENT_SECRET"),
    "scopes": [os.getenv("SCOPES")]
}

def get_token():
    app = ConfidentialClientApplication(
        config["client_id"], authority=config["authority"], client_credential=config["client_secret"]
    )
    result = app.acquire_token_for_client(scopes=config["scopes"])
    if "access_token" in result:
        return result["access_token"]
    else:
        raise Exception("Could not obtain access token")

def download_file(access_token, site_id, item_id, file_path):
    url = f"https://graph.microsoft.com/v1.0/sites/{site_id}/drive/items/{item_id}/content"
    headers = {
        "Authorization": f"Bearer {access_token}"
    }
    response = requests.get(url, headers=headers)
    response.raise_for_status()
    with open(file_path, "wb") as file:
        file.write(response.content)
    print(f"File downloaded to {file_path}")

def upload_file(access_token, site_id, item_id, file_path):
    url = f"https://graph.microsoft.com/v1.0/sites/{site_id}/drive/items/{item_id}/content"
    headers = {
        "Authorization": f"Bearer {access_token}",
        "Content-Type": "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    }
    with open(file_path, "rb") as file:
        response = requests.put(url, headers=headers, data=file)
    response.raise_for_status()
    print(f"File uploaded from {file_path} to item ID {item_id}")

def read_excel_data(file_path, sheet_name, start_row, end_row):
    workbook = load_workbook(filename=file_path, data_only=True)
    sheet = workbook[sheet_name]
    data = []
    for row in sheet.iter_rows(min_row=start_row, max_row=end_row, values_only=True):
        row_data = [cell if cell is not None else "" for cell in row]
        if any(row_data):
            data.append(row_data)
    print("Data read from Excel:")
    for row in data:
        print(row)
    return data

def align_data_with_headers(data, headers):
    aligned_data = []
    for row in data:
        if row[1] == "Totals":
            row = ["Totals"] + [""] * 2 + row[3:]  # Adjust for merged cells
        row = row[:len(headers)]
        if len(row) < len(headers):
            row += [""] * (len(headers) - len(row))
        aligned_data.append(row)
    print("Aligned data with headers:")
    for row in aligned_data:
        print(row)
    return aligned_data

def format_percentage(value):
    try:
        return f"{float(value):.2f}%"
    except (ValueError, TypeError):
        return str(value)

def update_powerpoint(ppt_path, data):
    presentation = Presentation(ppt_path)
    slide = presentation.slides[5]  # Slide 6 (0-indexed)

    table = None
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            break

    if not table:
        raise Exception("No table found on slide 6")

    headers = ["Brand", "Channel", "Following", "Reach", "% Reach MoM", "# New page likes & followers", "% follower change MoM", "# of Page Visits", "% Page Visits MoM"]

    # Align data with headers
    data = align_data_with_headers(data, headers)

    # Ensure table dimensions are sufficient
    if len(data) + 1 != len(table.rows) or len(headers) != len(table.columns):
        print(f"Data length: {len(data)}, Table rows: {len(table.rows)}, Headers: {len(headers)}, Table columns: {len(table.columns)}")  # Debug: Print dimensions
        raise Exception("Table dimensions do not match the data dimensions")

    # Clear existing table data
    for row in table.rows:
        for cell in row.cells:
            cell.text = ""

    # Fill headers
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.name = 'Arial'
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Fill in new data
    for i, row_data in enumerate(data):
        for j, cell_data in enumerate(row_data):
            cell = table.cell(i + 1, j)
            if j in [4, 6, 8]:  # Columns with percentage values
                cell.text = format_percentage(cell_data)
            else:
                cell.text = str(cell_data)
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].font.name = 'Arial'
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    presentation.save(ppt_path)
    print(f"Updated PowerPoint saved at {ppt_path}")

def main():
    if len(sys.argv) != 3:
        print("Usage: python update_ppt.py <sourceFileId> <destinationFileId>")
        sys.exit(1)

    source_file_id = sys.argv[1]
    destination_file_id = sys.argv[2]

    access_token = get_token()
    site_id = os.getenv("SITE_ID")
    
    excel_path = "source.xlsx"
    ppt_path = "destination.pptx"

    # Download files from SharePoint
    download_file(access_token, site_id, source_file_id, excel_path)

    # Read data from Excel
    excel_data = read_excel_data(excel_path, "For Monthly Reports", 13, 21)
    
    # Check if data includes "Totals" row
    print("Excel Data with Totals Row Included:")
    print(excel_data)

    # Update PowerPoint with Excel data
    update_powerpoint(ppt_path, excel_data)

    # Upload the updated PowerPoint back to SharePoint
    upload_file(access_token, site_id, destination_file_id, ppt_path)

if __name__ == "__main__":
    main()
